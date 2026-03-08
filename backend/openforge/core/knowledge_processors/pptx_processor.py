"""
PPTX processor for OpenForge Knowledge System.

Processes uploaded PowerPoint files through:
1. Text extraction using python-pptx
2. Metadata extraction
3. Thumbnail generation (placeholder)
4. Chunking and embedding
"""
import logging
from pathlib import Path
from uuid import UUID
from typing import Optional

from openforge.config import get_settings
from openforge.core.content_processors.base import ContentProcessor, ProcessorResult

logger = logging.getLogger("openforge.pptx_processor")


class PptxProcessor(ContentProcessor):
    """Process PowerPoint presentations for knowledge storage and retrieval."""

    name = "pptx"
    supported_types = [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ]
    supported_extensions = [".ppt", ".pptx"]

    def __init__(self):
        self.settings = get_settings()

    async def process(
        self,
        file_path: str,
        workspace_id: UUID,
        knowledge_id: Optional[UUID] = None,
        **kwargs,
    ) -> ProcessorResult:
        """
        Process a PowerPoint presentation.

        Args:
            file_path: Path to the PPTX file
            workspace_id: UUID of the workspace
            knowledge_id: Optional UUID of the knowledge entry
            **kwargs: Additional options

        Returns:
            ProcessorResult with extracted text and metadata
        """
        result = ProcessorResult(success=False)

        pptx_path = Path(file_path)
        if not pptx_path.exists():
            result.error = f"PowerPoint file not found: {file_path}"
            logger.error(f"PowerPoint file not found: {file_path}")
            return result

        try:
            # Extract text using python-pptx
            text_content, metadata = await self._extract_content(pptx_path)
            result.content = text_content
            result.extracted_text = text_content
            result.metadata = metadata

            # Generate thumbnail (placeholder for PowerPoint files)
            result.thumbnail_path = await self._generate_thumbnail(pptx_path, workspace_id, knowledge_id)

            result.success = True

            # Embed if knowledge_id provided
            if knowledge_id and text_content.strip():
                try:
                    from openforge.core.knowledge_processor import knowledge_processor
                    await knowledge_processor.process_knowledge(
                        knowledge_id=knowledge_id,
                        workspace_id=workspace_id,
                        content=text_content,
                        knowledge_type="pptx",
                        title=result.ai_title,
                        tags=[],
                    )
                    result.embedded = True
                except Exception as e:
                    logger.warning(f"Failed to embed PPTX content: {e}")

        except Exception as e:
            logger.exception(f"Error processing PPTX {knowledge_id}: {e}")
            result.error = str(e)

        return result

    async def _extract_content(self, pptx_path: Path) -> tuple[str, dict]:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            slides_text = []
            metadata = {"slide_count": len(prs.slides)}

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())

                if slide_parts:
                    slides_text.append(f"## Slide {slide_num}\n" + "\n".join(slide_parts))

            return "\n\n".join(slides_text), metadata
        except ImportError:
            logger.warning("python-pptx not installed, cannot extract PPTX content")
            return "", {}
        except Exception as e:
            logger.error(f"Failed to extract PPTX content: {e}")
            return "", {}

    async def _generate_thumbnail(
        self, pptx_path: Path, workspace_id: UUID, knowledge_id: Optional[UUID]
    ) -> Optional[str]:
        """Generate a placeholder thumbnail for PowerPoint files."""
        if not knowledge_id:
            return None

        try:
            from PIL import Image, ImageDraw

            workspace_dir = Path(self.settings.workspace_root) / str(workspace_id)
            thumbnail_dir = workspace_dir / "thumbnails"
            thumbnail_dir.mkdir(parents=True, exist_ok=True)

            thumbnail_path = thumbnail_dir / f"{knowledge_id}.webp"

            # Create a simple placeholder image
            img = Image.new("RGB", (200, 260), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # Draw PowerPoint-like icon
            draw.rectangle([20, 20, 180, 240], fill=(209, 69, 32), outline=(160, 50, 20), width=2)
            draw.text((70, 110), "P", fill=(255, 255, 255))
            draw.text((40, 150), "PPTX", fill=(255, 255, 255))

            img.save(thumbnail_path, "WEBP", quality=85)
            return f"thumbnails/{knowledge_id}.webp"
        except Exception as e:
            logger.warning(f"Failed to generate PPTX thumbnail: {e}")
            return None
