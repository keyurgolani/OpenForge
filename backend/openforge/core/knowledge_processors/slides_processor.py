"""Slides Processing Pipeline.

1. Read all slides via python-pptx (title + body text + speaker notes)
2. Extract metadata (slide count, has speaker notes)
3. Generate thumbnail via LibreOffice
4. Chunk text → embed → store in Qdrant
"""
from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Optional
from uuid import UUID

logger = logging.getLogger("openforge.processors.slides")


class SlidesProcessor:
    """Complete slides (PPTX) knowledge processing pipeline."""

    async def process(
        self,
        knowledge_id: UUID,
        file_path: str,
        workspace_id: UUID,
        db_session=None,
    ) -> dict:
        """Run the full slides processing pipeline. Returns metadata dict."""
        result = {
            "metadata": {},
            "text": "",
        }

        # ── Step 1: Extract text ──
        try:
            result["text"], result["metadata"] = self._extract_text_and_metadata(file_path)
        except Exception as e:
            logger.warning("Slides extraction failed for %s: %s", knowledge_id, e)

        # ── Step 2: Thumbnail ──
        thumbnail_path = None
        try:
            from openforge.config import get_settings
            from openforge.core.knowledge_processors.thumbnail_utils import generate_office_thumbnail

            settings = get_settings()
            thumbnails_dir = os.path.join(settings.uploads_root, "knowledge-thumbnails")
            os.makedirs(thumbnails_dir, exist_ok=True)
            thumb_file = os.path.join(thumbnails_dir, f"{knowledge_id}.webp")
            if generate_office_thumbnail(file_path, thumb_file):
                thumbnail_path = thumb_file
        except Exception as e:
            logger.warning("Slides thumbnail generation failed for %s: %s", knowledge_id, e)

        # ── Step 3: Embed text ──
        if result["text"] and len(result["text"].strip()) >= 20:
            try:
                await self._embed_text(knowledge_id, workspace_id, result["text"])
            except Exception as e:
                logger.warning("Slides text embedding failed for %s: %s", knowledge_id, e)

        metadata = result["metadata"]
        return {
            "thumbnail_path": thumbnail_path,
            "file_metadata": {
                "slide_count": metadata.get("slide_count"),
                "has_notes": metadata.get("has_notes"),
                "slide_titles": metadata.get("slide_titles", []),
            },
            "content": result["text"],
            "ai_title": metadata.get("presentation_title")
                or Path(file_path).stem.replace("_", " ").replace("-", " ").title(),
        }

    def _extract_text_and_metadata(self, file_path: str) -> tuple[str, dict]:
        """Extract text from all slides as markdown sections."""
        from pptx import Presentation

        prs = Presentation(file_path)
        parts = []
        slide_titles = []
        has_notes = False
        presentation_title = None

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_title = ""
            body_texts = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Check if this is a title shape
                if shape.shape_type is not None and hasattr(shape, "placeholder_format"):
                    ph = shape.placeholder_format
                    if ph is not None and ph.idx is not None and ph.idx in (0, 1):
                        # idx 0 = title, idx 1 = center title
                        slide_title = text
                        if slide_num == 1 and not presentation_title:
                            presentation_title = text
                        continue

                body_texts.append(text)

            # Extract speaker notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_notes = True

            # Build slide section
            title_display = slide_title or f"Slide {slide_num}"
            slide_titles.append(title_display)
            parts.append(f"## {title_display}\n")

            if body_texts:
                parts.append("\n".join(body_texts))

            if notes_text:
                parts.append(f"\n> **Speaker Notes:** {notes_text}")

            parts.append("")  # Blank line between slides

        full_text = "\n".join(parts)

        metadata = {
            "slide_count": len(prs.slides),
            "has_notes": has_notes,
            "slide_titles": slide_titles,
            "presentation_title": presentation_title,
        }

        return full_text[:100000], metadata

    async def _embed_text(
        self, knowledge_id: UUID, workspace_id: UUID, text: str
    ) -> None:
        """Embed extracted text into openforge_knowledge collection."""
        from openforge.core.knowledge_processor import knowledge_processor

        await knowledge_processor.process_knowledge(
            knowledge_id=knowledge_id,
            workspace_id=workspace_id,
            content=text,
            knowledge_type="slides",
            title=None,
            tags=[],
        )


slides_processor = SlidesProcessor()
