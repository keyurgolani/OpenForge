"""Text extraction slot backends for the pipeline framework.

Backends:
- MarkerBackend: PDF text extraction via Marker subprocess with PyMuPDF fallback
- DoclingBackend: PDF text+table extraction via IBM Docling subprocess
- DoclingTableBackend: Structured table extraction via Docling
- EmbeddedImageBackend: Embedded image extraction from documents via PyMuPDF
- PyMuPDFMetadataBackend: PDF metadata extraction via PyMuPDF
- DocxBackend: DOCX text extraction with heading structure preserved
- XlsxBackend: XLSX sheet extraction as markdown tables
- PptxBackend: PPTX slide text extraction with slide numbers
- PlainTextBackend: Passthrough for note/bookmark/gist/journal content
"""
from __future__ import annotations

import asyncio
import json
import logging
import os
import sys
import tempfile
from pathlib import Path

from openforge.core.pipeline.registry import register_backend
from openforge.core.pipeline.types import SlotContext, SlotOutput

logger = logging.getLogger(__name__)

# Maximum characters to extract from any single document.
MAX_TEXT_LENGTH = 100_000

# Timeout for the Marker subprocess (seconds).
MARKER_SUBPROCESS_TIMEOUT = 300

# ── Marker subprocess script ──
# Runs in a separate Python process so OOM/crash can't kill the API server.
_MARKER_SUBPROCESS_SCRIPT = """\
import sys
import json

file_path = sys.argv[1]
out_path = sys.argv[2]

try:
    from marker.converters.pdf import PdfConverter
    from marker.models import create_model_dict

    models = create_model_dict()
    converter = PdfConverter(artifact_dict=models)
    rendered = converter(file_path)
    text = rendered.markdown or ""

    with open(out_path, "w") as f:
        json.dump({"text": text[:100000]}, f)

except Exception as e:
    with open(out_path, "w") as f:
        json.dump({"text": "", "error": str(e)}, f)
    sys.exit(1)
"""


# ---------------------------------------------------------------------------
# MarkerBackend — PDF text extraction
# ---------------------------------------------------------------------------

class MarkerBackend:
    """PDF text extraction via Marker subprocess with PyMuPDF fallback."""

    slot_type = "text_extraction"
    backend_name = "marker"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            text = await self._extract_marker_subprocess(file_path)
            if not text:
                text = await asyncio.to_thread(self._extract_pymupdf_fallback, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            logger.warning("MarkerBackend failed for %s: %s", file_path, e)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    async def _extract_marker_subprocess(self, file_path: str) -> str:
        """Run Marker in an isolated subprocess."""
        out_path: str | None = None
        try:
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".json", delete=False
            ) as out_file:
                out_path = out_file.name

            env = {**os.environ}
            try:
                from openforge.config import get_settings
                marker_dir = str(Path(get_settings().models_root) / "marker")
                env["DATALAB_MODELS_DIR"] = marker_dir
            except Exception:
                pass

            proc = await asyncio.create_subprocess_exec(
                sys.executable, "-c", _MARKER_SUBPROCESS_SCRIPT, file_path, out_path,
                env=env,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )

            try:
                stdout, stderr = await asyncio.wait_for(
                    proc.communicate(), timeout=MARKER_SUBPROCESS_TIMEOUT
                )
            except asyncio.TimeoutError:
                proc.kill()
                await proc.communicate()
                logger.warning(
                    "Marker subprocess timed out after %ds, falling back to PyMuPDF",
                    MARKER_SUBPROCESS_TIMEOUT,
                )
                return ""

            if proc.returncode != 0:
                stderr_text = stderr.decode(errors="replace")[-500:] if stderr else ""
                logger.warning(
                    "Marker subprocess failed (exit %d): %s — falling back to PyMuPDF",
                    proc.returncode, stderr_text,
                )
                return ""

            if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                with open(out_path, "r") as f:
                    data = json.load(f)
                text = data.get("text", "")
                if text:
                    logger.info("Marker extraction succeeded (%d chars)", len(text))
                return text[:MAX_TEXT_LENGTH]
            return ""

        except Exception as e:
            logger.warning("Marker subprocess error: %s — falling back to PyMuPDF", e)
            return ""
        finally:
            if out_path:
                try:
                    if os.path.exists(out_path):
                        os.unlink(out_path)
                except Exception:
                    pass

    @staticmethod
    def _extract_pymupdf_fallback(file_path: str) -> str:
        """Extract text from all pages using PyMuPDF (basic fallback)."""
        import fitz

        doc = fitz.open(file_path)
        text_parts: list[str] = []

        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text("text")
            if text and text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{text.strip()}")

        doc.close()
        return "\n\n".join(text_parts)[:MAX_TEXT_LENGTH]


# ── Docling subprocess script ──
# Runs in a separate Python process so OOM/crash can't kill the API server.
_DOCLING_SUBPROCESS_SCRIPT = """\
import sys
import json

file_path = sys.argv[1]
out_path = sys.argv[2]

try:
    from docling.document_converter import DocumentConverter

    converter = DocumentConverter()
    result = converter.convert(file_path)
    doc = result.document

    # Full markdown export
    text = doc.export_to_markdown()

    # Extract tables separately as markdown
    tables = []
    for table in doc.tables:
        try:
            tables.append(table.export_to_markdown())
        except Exception:
            pass

    with open(out_path, "w") as f:
        json.dump({"text": text[:100000], "tables": tables}, f)

except Exception as e:
    with open(out_path, "w") as f:
        json.dump({"text": "", "tables": [], "error": str(e)}, f)
    sys.exit(1)
"""


# ---------------------------------------------------------------------------
# DoclingBackend — PDF text+table extraction via IBM Docling
# ---------------------------------------------------------------------------

class DoclingBackend:
    """Document text+table extraction via IBM Docling with fallback."""

    slot_type = "text_extraction"
    backend_name = "docling"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            text, tables = await self._extract_docling_subprocess(file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            metadata = {}
            if tables:
                metadata["tables"] = tables
                metadata["table_count"] = len(tables)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                metadata=metadata,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            logger.warning("DoclingBackend failed for %s: %s", file_path, e)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    async def _extract_docling_subprocess(self, file_path: str) -> tuple[str, list[str]]:
        """Run Docling in an isolated subprocess."""
        out_path: str | None = None
        try:
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".json", delete=False
            ) as out_file:
                out_path = out_file.name

            env = {**os.environ}
            try:
                from openforge.common.config import get_settings
                docling_dir = str(Path(get_settings().models_root) / "docling")
                env["HF_HOME"] = docling_dir
            except Exception:
                pass

            proc = await asyncio.create_subprocess_exec(
                sys.executable, "-c", _DOCLING_SUBPROCESS_SCRIPT, file_path, out_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                env=env,
            )
            try:
                stdout, stderr = await asyncio.wait_for(
                    proc.communicate(), timeout=MARKER_SUBPROCESS_TIMEOUT
                )
            except asyncio.TimeoutError:
                proc.kill()
                await proc.communicate()
                logger.warning(
                    "Docling subprocess timed out after %ds",
                    MARKER_SUBPROCESS_TIMEOUT,
                )
                return "", []

            if proc.returncode != 0:
                stderr_text = stderr.decode(errors="replace")[-500:] if stderr else ""
                logger.warning(
                    "Docling subprocess failed (exit %d): %s",
                    proc.returncode, stderr_text,
                )
                return "", []

            if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                with open(out_path, "r") as f:
                    data = json.load(f)
                text = data.get("text", "")
                tables = data.get("tables", [])
                if text:
                    logger.info(
                        "Docling extraction succeeded (%d chars, %d tables)",
                        len(text), len(tables),
                    )
                return text[:MAX_TEXT_LENGTH], tables
            return "", []

        except Exception as e:
            logger.warning("Docling subprocess error: %s", e)
            return "", []
        finally:
            if out_path:
                try:
                    if os.path.exists(out_path):
                        os.unlink(out_path)
                except Exception:
                    pass


# ---------------------------------------------------------------------------
# DoclingTableBackend — structured table extraction via Docling
# ---------------------------------------------------------------------------

class DoclingTableBackend:
    """Structured table extraction via Docling — returns tables as markdown."""

    slot_type = "table_extraction"
    backend_name = "docling"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            _, tables = await DoclingBackend()._extract_docling_subprocess(file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            text = "\n\n".join(tables) if tables else ""
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                metadata={"table_count": len(tables)},
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )


# ---------------------------------------------------------------------------
# EmbeddedImageBackend — extract embedded images from documents via PyMuPDF
# ---------------------------------------------------------------------------

class EmbeddedImageBackend:
    """Extract embedded images from documents using PyMuPDF."""

    slot_type = "embedded_image_extraction"
    backend_name = "pymupdf"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            result = await asyncio.to_thread(self._extract_images, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=result["text"],
                metadata={"image_count": result["count"], "image_paths": result["paths"]},
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    @staticmethod
    def _extract_images(file_path: str) -> dict:
        import fitz

        doc = fitz.open(file_path)
        base_dir = Path(file_path).parent / f"{Path(file_path).stem}_images"
        base_dir.mkdir(exist_ok=True)

        paths: list[str] = []
        for page_num in range(doc.page_count):
            for img_idx, img in enumerate(doc[page_num].get_images(full=True)):
                xref = img[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n > 4:  # CMYK
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_path = str(base_dir / f"page{page_num + 1}_img{img_idx + 1}.png")
                    pix.save(img_path)
                    paths.append(img_path)
                except Exception:
                    pass

        doc.close()
        text = f"Extracted {len(paths)} embedded images." if paths else "No embedded images found."
        return {"text": text, "count": len(paths), "paths": paths}


# ---------------------------------------------------------------------------
# PyMuPDFMetadataBackend — PDF metadata extraction
# ---------------------------------------------------------------------------

class PyMuPDFMetadataBackend:
    """Extract PDF metadata (title, author, dates, page count) via PyMuPDF."""

    slot_type = "metadata_extraction"
    backend_name = "pymupdf-meta"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            metadata = await asyncio.to_thread(self._extract, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            parts = []
            if metadata.get("title"):
                parts.append(f"Title: {metadata['title']}")
            if metadata.get("author"):
                parts.append(f"Author: {metadata['author']}")
            if metadata.get("page_count"):
                parts.append(f"Pages: {metadata['page_count']}")
            text = "\n".join(parts)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                metadata=metadata,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    @staticmethod
    def _extract(file_path: str) -> dict:
        import fitz

        doc = fitz.open(file_path)
        meta = doc.metadata or {}
        result = {
            "title": meta.get("title", ""),
            "author": meta.get("author", ""),
            "subject": meta.get("subject", ""),
            "creator": meta.get("creator", ""),
            "producer": meta.get("producer", ""),
            "page_count": doc.page_count,
        }
        doc.close()
        return result


# ---------------------------------------------------------------------------
# DocxBackend — DOCX text extraction
# ---------------------------------------------------------------------------

class DocxBackend:
    """DOCX text extraction with heading structure preserved."""

    slot_type = "text_extraction"
    backend_name = "python-docx"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            text, metadata = await asyncio.to_thread(self._extract, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                metadata=metadata,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            logger.warning("DocxBackend failed for %s: %s", file_path, e)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    @staticmethod
    def _extract(file_path: str) -> tuple[str, dict]:
        """Extract text from DOCX with heading structure preserved."""
        import docx

        doc = docx.Document(file_path)
        parts: list[str] = []

        heading_map = {
            "Heading 1": "# ",
            "Heading 2": "## ",
            "Heading 3": "### ",
            "Heading 4": "#### ",
            "Heading 5": "##### ",
            "Heading 6": "###### ",
        }

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""
            prefix = heading_map.get(style_name, "")

            if prefix:
                parts.append(f"\n{prefix}{text}\n")
            elif style_name.startswith("List"):
                parts.append(f"- {text}")
            else:
                parts.append(text)

        # Extract tables as markdown
        for table in doc.tables:
            table_rows: list[str] = []
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append("| " + " | ".join(cells) + " |")
                if row_idx == 0:
                    table_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if table_rows:
                parts.append("\n" + "\n".join(table_rows) + "\n")

        full_text = "\n".join(parts)[:MAX_TEXT_LENGTH]

        # Metadata
        core = doc.core_properties
        word_count = sum(len(p.text.split()) for p in doc.paragraphs if p.text.strip())
        metadata = {
            "author": core.author or "",
            "title": core.title or "",
            "word_count": word_count,
            "paragraph_count": sum(1 for p in doc.paragraphs if p.text.strip()),
            "section_count": len(doc.sections),
        }

        return full_text, metadata


# ---------------------------------------------------------------------------
# XlsxBackend — XLSX sheet extraction
# ---------------------------------------------------------------------------

class XlsxBackend:
    """XLSX sheet extraction as markdown tables."""

    slot_type = "text_extraction"
    backend_name = "openpyxl"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            text, metadata = await asyncio.to_thread(self._extract, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                metadata=metadata,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            logger.warning("XlsxBackend failed for %s: %s", file_path, e)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    @staticmethod
    def _extract(file_path: str) -> tuple[str, dict]:
        """Extract all sheets as markdown tables."""
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts: list[str] = []
        sheet_names: list[str] = []
        total_rows = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            sheet_names.append(sheet_name)

            if not rows:
                continue

            # Filter out completely empty rows
            non_empty_rows: list[list[str]] = []
            max_cols = 0
            for row in rows:
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    non_empty_rows.append(cells)
                    max_cols = max(max_cols, len(cells))

            if not non_empty_rows:
                continue

            total_rows += len(non_empty_rows)

            # Build markdown table
            parts.append(f"## {sheet_name}\n")

            # Limit to first 200 rows
            display_rows = non_empty_rows[:200]

            # Pad cells to uniform column count
            for row in display_rows:
                while len(row) < max_cols:
                    row.append("")

            # Header row
            header = display_rows[0]
            parts.append("| " + " | ".join(header) + " |")
            parts.append("| " + " | ".join(["---"] * max_cols) + " |")

            # Data rows
            for row in display_rows[1:]:
                parts.append("| " + " | ".join(row) + " |")

            if len(non_empty_rows) > 200:
                parts.append(f"\n*... {len(non_empty_rows) - 200} more rows truncated*\n")

            parts.append("")  # Blank line between sheets

        wb.close()

        full_text = "\n".join(parts)[:MAX_TEXT_LENGTH]
        metadata = {
            "sheet_names": sheet_names,
            "total_rows": total_rows,
            "total_sheets": len(sheet_names),
        }

        return full_text, metadata


# ---------------------------------------------------------------------------
# PptxBackend — PPTX slide text extraction
# ---------------------------------------------------------------------------

class PptxBackend:
    """PPTX slide text extraction with slide numbers."""

    slot_type = "text_extraction"
    backend_name = "python-pptx"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            text, metadata = await asyncio.to_thread(self._extract, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                metadata=metadata,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            logger.warning("PptxBackend failed for %s: %s", file_path, e)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    @staticmethod
    def _extract(file_path: str) -> tuple[str, dict]:
        """Extract text from all slides as markdown sections."""
        from pptx import Presentation

        prs = Presentation(file_path)
        parts: list[str] = []
        slide_titles: list[str] = []
        has_notes = False
        presentation_title: str | None = None

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_title = ""
            body_texts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Check if this is a title shape
                if shape.shape_type is not None and hasattr(shape, "placeholder_format"):
                    ph = shape.placeholder_format
                    if ph is not None and ph.idx is not None and ph.idx in (0, 1):
                        slide_title = text
                        if slide_num == 1 and not presentation_title:
                            presentation_title = text
                        continue

                body_texts.append(text)

            # Extract speaker notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_notes = True

            # Build slide section
            title_display = slide_title or f"Slide {slide_num}"
            slide_titles.append(title_display)
            parts.append(f"## {title_display}\n")

            if body_texts:
                parts.append("\n".join(body_texts))

            if notes_text:
                parts.append(f"\n> **Speaker Notes:** {notes_text}")

            parts.append("")  # Blank line between slides

        full_text = "\n".join(parts)[:MAX_TEXT_LENGTH]

        metadata = {
            "slide_count": len(prs.slides),
            "has_notes": has_notes,
            "slide_titles": slide_titles,
            "presentation_title": presentation_title,
        }

        return full_text, metadata


# ---------------------------------------------------------------------------
# PlainTextBackend — passthrough for text-based content
# ---------------------------------------------------------------------------

class PlainTextBackend:
    """Passthrough backend for note/bookmark/gist/journal content."""

    slot_type = "text_extraction"
    backend_name = "plaintext"

    async def run(self, file_path: str, context: SlotContext) -> SlotOutput:
        import time

        start = time.monotonic()
        try:
            text = await asyncio.to_thread(self._read_file, file_path)
            elapsed = int((time.monotonic() - start) * 1000)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                text=text,
                duration_ms=elapsed,
            )
        except Exception as e:
            elapsed = int((time.monotonic() - start) * 1000)
            logger.warning("PlainTextBackend failed for %s: %s", file_path, e)
            return SlotOutput(
                slot_type=self.slot_type,
                backend_name=self.backend_name,
                success=False,
                error=str(e),
                duration_ms=elapsed,
            )

    @staticmethod
    def _read_file(file_path: str) -> str:
        """Read file content as UTF-8 text."""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()[:MAX_TEXT_LENGTH]


# ---------------------------------------------------------------------------
# Register all backends
# ---------------------------------------------------------------------------

register_backend("text_extraction", "marker", MarkerBackend())
register_backend("text_extraction", "docling", DoclingBackend())
register_backend("text_extraction", "python-docx", DocxBackend())
register_backend("text_extraction", "openpyxl", XlsxBackend())
register_backend("text_extraction", "python-pptx", PptxBackend())
register_backend("text_extraction", "plaintext", PlainTextBackend())
register_backend("table_extraction", "docling", DoclingTableBackend())
register_backend("embedded_image_extraction", "pymupdf", EmbeddedImageBackend())
register_backend("metadata_extraction", "pymupdf-meta", PyMuPDFMetadataBackend())
