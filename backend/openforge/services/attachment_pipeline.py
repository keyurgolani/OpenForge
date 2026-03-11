from __future__ import annotations

import re
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional

TEXT_FILE_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".csv",
    ".xml",
    ".yaml",
    ".yml",
}

URL_TRAILING_PUNCTUATION = ".,!?;:)]}\"'"
HTTP_URL_PATTERN = re.compile(r"https?://[^\s<>]+")


# ---------------------------------------------------------------------------
# Extractor registry
# ---------------------------------------------------------------------------

class AttachmentExtractor(ABC):
    """Base class for content extractors. Each pipeline type has one concrete extractor."""

    pipeline: str  # e.g. "text", "pdf", "image", "audio"

    @abstractmethod
    def matches(self, content_type: str, extension: str) -> bool:
        """Return True if this extractor handles the given content_type / file extension."""

    @abstractmethod
    async def extract(self, file_path: str) -> Optional[str]:
        """Extract and return text content, or None if nothing could be extracted."""


class TextAttachmentExtractor(AttachmentExtractor):
    """Extracts plain text from text-based file formats."""

    pipeline = "text"

    def matches(self, content_type: str, extension: str) -> bool:
        return content_type.startswith("text/") or extension in TEXT_FILE_EXTENSIONS

    async def extract(self, file_path: str) -> Optional[str]:
        import aiofiles
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = await f.read()
                raw = content[:50000]
            return (raw or "").strip() or None
        except Exception:
            return None


class PDFAttachmentExtractor(AttachmentExtractor):
    """Extracts text content from PDF files using PyMuPDF."""

    pipeline = "pdf"

    def matches(self, content_type: str, extension: str) -> bool:
        return content_type == "application/pdf" or extension == ".pdf"

    async def extract(self, file_path: str) -> Optional[str]:
        try:
            import fitz
            doc = fitz.open(file_path)
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            doc.close()
            text = "\n\n".join(text_parts)[:50000]
            return text.strip() or None
        except Exception:
            return None


class ImageAttachmentExtractor(AttachmentExtractor):
    """Extracts text from images using Tesseract OCR."""

    pipeline = "image"

    def matches(self, content_type: str, extension: str) -> bool:
        return content_type.startswith("image/")

    async def extract(self, file_path: str) -> Optional[str]:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return (text or "").strip()[:10000] or None
        except Exception:
            return None


class AudioAttachmentExtractor(AttachmentExtractor):
    """Extracts text from audio files using Whisper transcription."""

    pipeline = "audio"

    def matches(self, content_type: str, extension: str) -> bool:
        return content_type.startswith("audio/") or extension in {
            ".mp3", ".wav", ".ogg", ".flac", ".m4a", ".weba",
        }

    async def extract(self, file_path: str) -> Optional[str]:
        try:
            from openforge.core.knowledge_processors.audio_processor import (
                _get_whisper_model, _get_whisper_download_root,
            )
            import asyncio

            download_root = _get_whisper_download_root()
            # Use base model for chat attachments — fast and lightweight
            model = _get_whisper_model("base", download_root=download_root)
            result = await asyncio.to_thread(
                model.transcribe, file_path, fp16=False, verbose=False
            )
            text = result.get("text", "").strip()
            return text[:50000] or None
        except Exception:
            # Fall back to metadata-only
            try:
                import mutagen
                audio = mutagen.File(file_path)
                if audio and audio.info:
                    duration = getattr(audio.info, "length", 0)
                    minutes = int(duration // 60)
                    seconds = int(duration % 60)
                    return f"[Audio file: {minutes}:{seconds:02d} duration]"
            except Exception:
                pass
            return None


class DocumentAttachmentExtractor(AttachmentExtractor):
    """Extracts text from DOCX files using python-docx."""

    pipeline = "document"

    def matches(self, content_type: str, extension: str) -> bool:
        return (
            content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or content_type == "application/msword"
            or extension in {".docx", ".doc"}
        )

    async def extract(self, file_path: str) -> Optional[str]:
        try:
            import docx
            doc = docx.Document(file_path)
            parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
            text = "\n".join(parts)
            return text[:50000].strip() or None
        except Exception:
            return None


class SheetAttachmentExtractor(AttachmentExtractor):
    """Extracts text from XLSX/XLS files using openpyxl."""

    pipeline = "sheet"

    def matches(self, content_type: str, extension: str) -> bool:
        return (
            content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            or content_type == "application/vnd.ms-excel"
            or extension in {".xlsx", ".xls"}
        )

    async def extract(self, file_path: str) -> Optional[str]:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"Sheet: {sheet_name}")
                for row in ws.iter_rows(values_only=True, max_row=200):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(c.strip() for c in cells):
                        parts.append(" | ".join(cells))
            wb.close()
            text = "\n".join(parts)
            return text[:50000].strip() or None
        except Exception:
            return None


class SlidesAttachmentExtractor(AttachmentExtractor):
    """Extracts text from PPTX files using python-pptx."""

    pipeline = "slides"

    def matches(self, content_type: str, extension: str) -> bool:
        return (
            content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or content_type == "application/vnd.ms-powerpoint"
            or extension in {".pptx", ".ppt"}
        )

    async def extract(self, file_path: str) -> Optional[str]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_texts.append(text)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_texts.append(f"[Notes: {notes}]")
                if slide_texts:
                    parts.append(f"Slide {slide_num}: " + " | ".join(slide_texts))
            text = "\n".join(parts)
            return text[:50000].strip() or None
        except Exception:
            return None


_EXTRACTORS: list[AttachmentExtractor] = [
    TextAttachmentExtractor(),
    PDFAttachmentExtractor(),
    ImageAttachmentExtractor(),
    AudioAttachmentExtractor(),
    DocumentAttachmentExtractor(),
    SheetAttachmentExtractor(),
    SlidesAttachmentExtractor(),
]


def get_extractor(content_type: str | None, filename: str | None) -> Optional[AttachmentExtractor]:
    """Return the matching extractor for this file, or None if none is registered yet."""
    ct = (content_type or "").strip().lower()
    ext = Path((filename or "").strip()).suffix.lower()
    for extractor in _EXTRACTORS:
        if extractor.matches(ct, ext):
            return extractor
    return None


def resolve_attachment_pipeline(content_type: str | None, filename: str | None) -> str:
    """Return the pipeline identifier for this attachment."""
    extractor = get_extractor(content_type, filename)
    if extractor is not None:
        return extractor.pipeline
    return "deferred"


def extract_http_urls(text: str) -> list[str]:
    ordered_unique: list[str] = []
    seen: set[str] = set()

    for match in HTTP_URL_PATTERN.findall(text or ""):
        cleaned = match.rstrip(URL_TRAILING_PUNCTUATION)
        if cleaned and cleaned not in seen:
            seen.add(cleaned)
            ordered_unique.append(cleaned)

    return ordered_unique
